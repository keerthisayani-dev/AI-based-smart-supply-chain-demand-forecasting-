from pathlib import Path
from copy import deepcopy

import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(r"c:/Users/Keerthi Sayani/OneDrive/Documents/supply chain project/project_implementation data")
DATA_PATH = BASE_DIR / "Dataset" / "retail_store_inventory.csv"
REF_PPT = Path(r"C:/Users/Keerthi Sayani/OneDrive/Desktop/Deep Fake Analysis.pptx")
OUT_PPT = BASE_DIR / "Supply_Chain_Project_Presentation.pptx"
ASSET_DIR = BASE_DIR / "ppt_assets"
ASSET_DIR.mkdir(exist_ok=True)


def delete_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = slide_ids[i].rId
        prs.part.drop_rel(r_id)
        del slide_ids[i]


def style_title(shape, size=40):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)


def style_body(shape, size=20):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor(255, 255, 255)


def build_charts(df: pd.DataFrame):
    df["date"] = pd.to_datetime(df["date"], dayfirst=True, errors="coerce")

    # 1) Monthly demand trend
    monthly = df.groupby(df["date"].dt.to_period("M"))["units sold"].sum().to_timestamp()
    plt.figure(figsize=(10, 4.5))
    plt.plot(monthly.index, monthly.values, color="#00c2ff", linewidth=2.2)
    plt.title("Monthly Units Sold Trend (2022-2024)", fontsize=14)
    plt.xlabel("Month")
    plt.ylabel("Units Sold")
    plt.grid(alpha=0.25)
    plt.tight_layout()
    trend_path = ASSET_DIR / "monthly_trend.png"
    plt.savefig(trend_path, dpi=200)
    plt.close()

    # 2) Category contribution
    cat = df.groupby("category")["units sold"].sum().sort_values(ascending=False).head(10)
    plt.figure(figsize=(10, 5.2))
    plt.barh(cat.index[::-1], cat.values[::-1], color="#11b67a")
    plt.title("Top 10 Categories by Units Sold", fontsize=14)
    plt.xlabel("Total Units Sold")
    plt.tight_layout()
    cat_path = ASSET_DIR / "category_bar.png"
    plt.savefig(cat_path, dpi=200)
    plt.close()

    # 3) Location split pie
    loc = df.groupby("location")["units sold"].sum().sort_values(ascending=False)
    plt.figure(figsize=(7, 7))
    colors = ["#1f77b4", "#ff7f0e", "#2ca02c", "#d62728"]
    plt.pie(loc.values, labels=loc.index, autopct="%1.1f%%", startangle=140, colors=colors)
    plt.title("Regional Sales Share")
    plt.tight_layout()
    loc_path = ASSET_DIR / "location_pie.png"
    plt.savefig(loc_path, dpi=200)
    plt.close()

    # 4) Architecture flow diagram
    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.axis("off")

    nodes = [
        (0.04, 0.33, 0.2, 0.34, "Data Sources\n(Sales + Inventory CSV)"),
        (0.30, 0.33, 0.2, 0.34, "Preprocessing\n+ Feature Engineering"),
        (0.56, 0.33, 0.2, 0.34, "Forecast Engine\n(LinearRegression)"),
        (0.82, 0.33, 0.14, 0.34, "FastAPI +\nDashboard"),
    ]

    for x, y, w, h, label in nodes:
        box = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02", linewidth=2, edgecolor="#1b3c73", facecolor="#d7e6ff")
        ax.add_patch(box)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=10, weight="bold")

    for i in range(len(nodes) - 1):
        x1 = nodes[i][0] + nodes[i][2]
        x2 = nodes[i + 1][0]
        y = nodes[i][1] + nodes[i][3] / 2
        ax.annotate("", xy=(x2 - 0.01, y), xytext=(x1 + 0.01, y), arrowprops=dict(arrowstyle="->", lw=2, color="#1b3c73"))

    arch_path = ASSET_DIR / "architecture.png"
    plt.tight_layout()
    plt.savefig(arch_path, dpi=200)
    plt.close()

    return trend_path, cat_path, loc_path, arch_path


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Retail Inventory Forecasting System"
    subtitle = slide.placeholders[1]
    subtitle.text = "AI-powered demand forecasting for supply chain decisions\nProject presentation"
    style_title(slide.shapes.title, 42)
    style_body(subtitle, 20)


def add_content_slide(prs, title, bullets, image_path=None, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 30)

    tx_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.2), Inches(5.8), Inches(5.4))
    tf = tx_box.text_frame
    tf.word_wrap = True

    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(245, 245, 245)
        first = False

    if image_path:
        slide.shapes.add_picture(str(image_path), Inches(6.45), Inches(1.45), width=Inches(6.7), height=Inches(3.9))
        if caption:
            cap = slide.shapes.add_textbox(Inches(6.45), Inches(5.45), Inches(6.7), Inches(0.6))
            cap_tf = cap.text_frame
            cap_tf.text = caption
            cap_tf.paragraphs[0].font.size = Pt(14)
            cap_tf.paragraphs[0].font.color.rgb = RGBColor(240, 240, 240)


def main():
    df = pd.read_csv(DATA_PATH)
    trend_path, cat_path, loc_path, arch_path = build_charts(df)

    total_rows = len(df)
    date_min = pd.to_datetime(df["date"], dayfirst=True).min().date()
    date_max = pd.to_datetime(df["date"], dayfirst=True).max().date()
    total_units = df["units sold"].sum()
    avg_daily_units = df.groupby(pd.to_datetime(df["date"], dayfirst=True))["units sold"].sum().mean()
    avg_inventory = df["inventory level"].mean()
    low_inv_pct = (df["inventory level"] < 80).mean() * 100
    mae_forecast = (df["demand forecast"] - df["units sold"]).abs().mean()

    prs = Presentation(str(REF_PPT))
    delete_all_slides(prs)

    add_title_slide(prs)

    add_content_slide(
        prs,
        "Project Overview",
        [
            "Objective: build a forecasting-driven system to improve retail inventory planning and reduce stock imbalances.",
            "Scope: 5 stores, 30 products, 15 categories, and multivariate demand signals including price, weather, and promotions.",
            "Delivery: end-to-end pipeline with data ingestion, feature engineering, model training, API services, and dashboard outputs.",
            "Outcome: daily demand forecasts that support category-level replenishment and operational decision-making.",
        ],
    )

    add_content_slide(
        prs,
        "Dataset Profile",
        [
            f"Records: {total_rows:,} rows with 15 attributes from {date_min} to {date_max}.",
            "Core variables: inventory level, units sold, units ordered, demand forecast, price, discount, seasonality, and promotion flag.",
            f"Sales volume: {total_units:,.0f} total units sold, with an average of {avg_daily_units:,.0f} units/day.",
            f"Inventory status: mean inventory {avg_inventory:.1f}; low-stock observations (<80 units) are {low_inv_pct:.1f}% of all records.",
        ],
        image_path=trend_path,
        caption="Figure: Monthly demand trend used to identify seasonality and long-term movement.",
    )

    add_content_slide(
        prs,
        "System Architecture",
        [
            "Data layer loads CSV uploads and standardizes columns, dates, and category identifiers.",
            "Feature layer creates lag, rolling, and seasonal calendar signals for each category-level daily series.",
            "Model layer uses LinearRegression to learn demand behavior and generate multi-day forecasts.",
            "Service layer (FastAPI) exposes forecast endpoints and powers a role-based dashboard for business users.",
        ],
        image_path=arch_path,
        caption="Figure: End-to-end architecture from raw data to dashboard consumption.",
    )

    add_content_slide(
        prs,
        "Exploratory Insights",
        [
            "Category demand is concentrated: Toys, Electronics, and Furniture contribute the largest share of unit sales.",
            "Regional differences are visible, with East and West driving the majority of demand volume.",
            "Price, discount, and holiday/promotion variables exhibit measurable influence across categories.",
            "These patterns justify feature-rich forecasting rather than simple moving-average baselines.",
        ],
        image_path=cat_path,
        caption="Figure: Top category contribution highlights demand concentration by product group.",
    )

    add_content_slide(
        prs,
        "Regional Demand Distribution",
        [
            "Location-level sales share indicates uneven demand distribution across operating regions.",
            "This supports location-sensitive replenishment rather than a uniform inventory policy.",
            "Forecast outputs can be combined with warehouse allocation constraints for region-aware planning.",
            "The same framework can be extended for store-level service-level targets.",
        ],
        image_path=loc_path,
        caption="Figure: Regional share of units sold across North, South, East, and West.",
    )

    add_content_slide(
        prs,
        "Forecasting Methodology",
        [
            "Model: LinearRegression with lag-1/7/14/28 features, rolling statistics, and cyclical calendar encodings.",
            "Adaptive controls: holiday-aware adjustments and exogenous indicators such as price and promotions.",
            f"Baseline quality indicator: mean absolute gap between historical demand forecast and actual sales is {mae_forecast:.2f} units.",
            "Operational use: generate next-day to next-week category forecasts for replenishment and procurement planning.",
        ],
    )

    add_content_slide(
        prs,
        "Business Impact",
        [
            "Improves replenishment timing by giving planners forward visibility into likely demand surges and dips.",
            "Reduces stockout risk and overstock by aligning purchase quantities with predicted consumption patterns.",
            "Creates a single decision workflow that combines forecasting outputs with inventory-health indicators.",
            "Supports measurable KPIs: service level, stockout rate, holding cost, and forecast error trend over time.",
        ],
    )

    add_content_slide(
        prs,
        "Implementation Challenges",
        [
            "Data quality variation across uploads required robust date parsing and column-alias normalization.",
            "Irregular observations demanded daily resampling and missing-value handling before model fitting.",
            "Model governance requires continuous monitoring to detect drift across categories and seasons.",
            "Role-based access and audit history were integrated to improve operational reliability and accountability.",
        ],
    )

    add_content_slide(
        prs,
        "Conclusion & Next Steps",
        [
            "The project demonstrates a complete demand-forecasting workflow from raw data to actionable dashboard insights.",
            "Current system is production-ready for category-level planning and can scale to store-SKU granularity.",
            "Next steps: add probabilistic intervals, automated retraining, and optimization-based reorder recommendations.",
            "Result: a practical AI-enabled supply chain decision support platform for retail operations.",
        ],
    )

    prs.save(str(OUT_PPT))
    print(f"Created: {OUT_PPT}")


if __name__ == "__main__":
    main()
