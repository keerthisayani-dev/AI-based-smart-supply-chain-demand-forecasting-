from pathlib import Path
import re

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(r"c:/Users/Keerthi Sayani/OneDrive/Documents/supply chain project/project_implementation data")
DEEPFAKE_TEMPLATE = Path(r"C:/Users/Keerthi Sayani/OneDrive/Desktop/Deep Fake Analysis.pptx")
REVIEW_TEMPLATE = Path(r"C:/Users/Keerthi Sayani/Downloads/Review_Template(1).pptx")
DATA_PATH = BASE_DIR / "Dataset" / "retail_store_inventory.csv"
OUT_PPT = BASE_DIR / "Backend" / "Supply_Chain_Review_Deck_From_Template.pptx"


def delete_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = slide_ids[i].rId
        prs.part.drop_rel(r_id)
        del slide_ids[i]


def extract_agenda_headings(review_ppt: Path):
    p = Presentation(str(review_ppt))
    for s in p.slides:
        for sh in s.shapes:
            if hasattr(sh, "text"):
                txt = (sh.text or "").strip()
                if "AGENDA" in txt.upper() and len(txt) > 10:
                    body = re.sub(r"(?i)AGENDA", "", txt).strip("\n :")
                    parts = [x.strip() for x in re.split(r"\||\n|•", body) if x.strip()]
                    if parts:
                        return parts
    return [
        "Problem Statement",
        "Objective",
        "Methodology",
        "System Architecture",
        "Technology Stack",
        "Implementation and Present Results",
        "Future Work",
    ]


def style_title(shape, size=32):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor(255, 255, 255)


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 30)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    for line in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(242, 246, 255)
        p.space_after = Pt(10)
        first = False


def main():
    agenda = extract_agenda_headings(REVIEW_TEMPLATE)

    df = pd.read_csv(DATA_PATH)
    df["date"] = pd.to_datetime(df["date"], dayfirst=True, errors="coerce")

    total_rows = len(df)
    date_min = df["date"].min().date()
    date_max = df["date"].max().date()
    stores = df["store id"].nunique()
    categories = df["category"].nunique()
    products = df["product id"].nunique()
    total_units = df["units sold"].sum()
    avg_inventory = df["inventory level"].mean()
    mae_gap = (df["demand forecast"] - df["units sold"]).abs().mean()

    content = {
        "Problem Statement": [
            "Retail inventory decisions are often reactive, causing stockouts and excess holding cost.",
            "Demand fluctuates by category, location, price, weather, and seasonal events.",
            "The project needs a reliable forecasting workflow to support daily replenishment planning.",
        ],
        "Objective": [
            "Build an AI-assisted demand forecasting system for supply-chain and inventory teams.",
            "Provide category-level forecasts through a web dashboard with role-based access control.",
            "Improve planning accuracy and reduce inventory imbalance across stores and regions.",
        ],
        "Methodology": [
            "Data preprocessing normalizes schema, parses mixed date formats, and handles missing values.",
            "Feature engineering adds lag, rolling, and calendar-seasonality signals.",
            "LinearRegression model is trained on category-level daily time series for next-day/multi-day prediction.",
        ],
        "System Architecture": [
            "Data Layer: CSV dataset upload and validation pipeline.",
            "Model Layer: forecasting_core with feature building and forecast generation.",
            "Service Layer: FastAPI endpoints for forecast, history, auth, and admin operations.",
            "Presentation Layer: HTML/CSS/JS dashboard for forecast visualization and export.",
        ],
        "Technology Stack": [
            "Backend: Python, FastAPI, Pandas, scikit-learn, SQLite.",
            "Frontend: HTML, CSS, JavaScript with Plotly-based visualization.",
            "Data: retail_store_inventory.csv (multi-factor retail demand dataset).",
            "Tooling: role-based authentication, upload management, and forecast history tracking.",
        ],
        "Implementation and Present Results": [
            f"Dataset processed: {total_rows:,} records from {date_min} to {date_max}.",
            f"Coverage: {stores} stores, {products} products, {categories} categories.",
            f"Total units sold in dataset: {total_units:,.0f}; average inventory level: {avg_inventory:.1f}.",
            f"Observed mean absolute gap (demand forecast vs units sold): {mae_gap:.2f} units.",
        ],
        "Future Work": [
            "Add probabilistic forecasting with confidence bands for risk-aware planning.",
            "Automate retraining and model monitoring for drift detection.",
            "Extend to SKU-store level optimization with reorder recommendation logic.",
            "Integrate procurement lead-time and supplier constraints for end-to-end planning.",
        ],
    }

    prs = Presentation(str(DEEPFAKE_TEMPLATE))
    delete_all_slides(prs)

    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "SUPPLY CHAIN FORECASTING SYSTEM"
    title.placeholders[1].text = "Review Presentation Based on Approved Agenda"
    style_title(title.shapes.title, 40)

    agenda_slide = prs.slides.add_slide(prs.slide_layouts[1])
    agenda_slide.shapes.title.text = "AGENDA"
    style_title(agenda_slide.shapes.title, 34)
    body = agenda_slide.placeholders[1].text_frame
    body.clear()
    for i, h in enumerate(agenda, 1):
        p = body.paragraphs[0] if i == 1 else body.add_paragraph()
        p.text = f"{i}. {h}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(245, 245, 245)

    for h in agenda:
        bullets = content.get(h, [
            "Section content aligned to project implementation and current progress.",
            "Detailed explanation can be expanded with additional metrics as needed.",
        ])
        add_bullets_slide(prs, h, bullets)

    thanks = prs.slides.add_slide(prs.slide_layouts[2])
    thanks.shapes.title.text = "Thank You"
    style_title(thanks.shapes.title, 44)

    prs.save(str(OUT_PPT))
    print(f"Created: {OUT_PPT}")
    print("Agenda used:", " | ".join(agenda))


if __name__ == "__main__":
    main()
