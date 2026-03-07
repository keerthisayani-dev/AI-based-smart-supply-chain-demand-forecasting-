from pathlib import Path

import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(r"c:/Users/Keerthi Sayani/OneDrive/Documents/supply chain project/project_implementation data")
BACKEND_DIR = BASE_DIR / "Backend"
DATA_PATH = BASE_DIR / "Dataset" / "retail_store_inventory.csv"
OUT_PPT = BACKEND_DIR / "Supply_Chain_Project_Detailed_Explanation_Till_Now.pptx"
ASSET_DIR = BACKEND_DIR / "ppt_assets"
ASSET_DIR.mkdir(exist_ok=True)


def build_charts(df: pd.DataFrame):
    df = df.copy()
    df["date"] = pd.to_datetime(df["date"], dayfirst=True, errors="coerce")

    monthly = df.groupby(df["date"].dt.to_period("M"))["units sold"].sum().to_timestamp()
    plt.figure(figsize=(10, 4.6))
    plt.plot(monthly.index, monthly.values, color="#1f77b4", linewidth=2.2)
    plt.title("Monthly Units Sold Trend (2022-2024)", fontsize=14)
    plt.xlabel("Month")
    plt.ylabel("Units Sold")
    plt.grid(alpha=0.25)
    plt.tight_layout()
    monthly_path = ASSET_DIR / "monthly_units_trend.png"
    plt.savefig(monthly_path, dpi=200)
    plt.close()

    category = df.groupby("category")["units sold"].sum().sort_values(ascending=False).head(10)
    plt.figure(figsize=(10, 5.2))
    plt.barh(category.index[::-1], category.values[::-1], color="#2ca02c")
    plt.title("Top 10 Categories by Total Units Sold", fontsize=14)
    plt.xlabel("Units Sold")
    plt.tight_layout()
    category_path = ASSET_DIR / "top_categories.png"
    plt.savefig(category_path, dpi=200)
    plt.close()

    loc = df.groupby("location")["units sold"].sum().sort_values(ascending=False)
    plt.figure(figsize=(7, 7))
    colors = ["#1f77b4", "#ff7f0e", "#2ca02c", "#d62728"]
    plt.pie(loc.values, labels=loc.index, autopct="%1.1f%%", startangle=145, colors=colors)
    plt.title("Sales Share by Region")
    plt.tight_layout()
    loc_path = ASSET_DIR / "regional_share.png"
    plt.savefig(loc_path, dpi=200)
    plt.close()

    fig, ax = plt.subplots(figsize=(11.5, 4.8))
    ax.axis("off")
    blocks = [
        (0.03, 0.33, 0.2, 0.35, "Dataset\nCSV Uploads"),
        (0.27, 0.33, 0.2, 0.35, "Preprocessing\n+ Validation"),
        (0.51, 0.33, 0.2, 0.35, "Forecasting Core\nFeature + Model"),
        (0.75, 0.33, 0.2, 0.35, "FastAPI +\nDashboard UI"),
    ]

    for x, y, w, h, label in blocks:
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02", linewidth=2, edgecolor="#244678", facecolor="#d9e7ff")
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=10, weight="bold")

    for i in range(len(blocks) - 1):
        x1 = blocks[i][0] + blocks[i][2]
        x2 = blocks[i + 1][0]
        y = blocks[i][1] + blocks[i][3] / 2
        ax.annotate("", xy=(x2 - 0.01, y), xytext=(x1 + 0.01, y), arrowprops=dict(arrowstyle="->", lw=2, color="#244678"))

    arch_path = ASSET_DIR / "system_architecture.png"
    plt.tight_layout()
    plt.savefig(arch_path, dpi=200)
    plt.close()

    return monthly_path, category_path, loc_path, arch_path


def style_title(title_shape, size=32):
    tf = title_shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(15, 35, 75)


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Supply Chain Forecasting Project"
    sub = slide.placeholders[1]
    sub.text = "Detailed explanation of implementation status till now"


def add_bullets_slide(prs: Presentation, title: str, bullets, image_path=None, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    style_title(slide.shapes.title, 30)

    txt = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(6.0), Inches(5.4))
    tf = txt.text_frame
    tf.word_wrap = True

    first = True
    for line in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(19)
        p.font.color.rgb = RGBColor(35, 45, 60)
        p.space_after = Pt(7)
        first = False

    if image_path:
        slide.shapes.add_picture(str(image_path), Inches(6.65), Inches(1.45), width=Inches(6.5), height=Inches(3.9))
        if caption:
            cap = slide.shapes.add_textbox(Inches(6.65), Inches(5.45), Inches(6.5), Inches(0.6))
            cap_tf = cap.text_frame
            cap_tf.text = caption
            cap_tf.paragraphs[0].font.size = Pt(13)
            cap_tf.paragraphs[0].font.color.rgb = RGBColor(70, 80, 95)


def main():
    df = pd.read_csv(DATA_PATH)
    monthly_path, category_path, loc_path, arch_path = build_charts(df)

    df["date"] = pd.to_datetime(df["date"], dayfirst=True, errors="coerce")
    total_rows = len(df)
    total_cols = len(df.columns)
    date_min = df["date"].min().date()
    date_max = df["date"].max().date()
    stores = df["store id"].nunique()
    products = df["product id"].nunique()
    categories = df["category"].nunique()
    total_units = df["units sold"].sum()
    avg_daily_units = df.groupby("date")["units sold"].sum().mean()
    avg_inventory = df["inventory level"].mean()
    low_stock_pct = (df["inventory level"] < 80).mean() * 100
    abs_gap = (df["demand forecast"] - df["units sold"]).abs().mean()

    prs = Presentation()

    add_title_slide(prs)

    add_bullets_slide(
        prs,
        "1. Project Objective",
        [
            "Goal: build an AI-enabled demand forecasting platform for retail supply-chain planning.",
            "Focus: predict short-term demand and support inventory decisions to reduce stockouts and overstock.",
            "Scope: complete stack delivery including data pipeline, model logic, API services, and business dashboard.",
            "Status: working end-to-end prototype with authentication, role controls, and downloadable outputs.",
        ],
    )

    add_bullets_slide(
        prs,
        "2. Dataset Summary",
        [
            f"Dataset size: {total_rows:,} rows x {total_cols} columns; period {date_min} to {date_max}.",
            f"Coverage: {stores} stores, {products} products, {categories} categories, and 4 regions.",
            "Signals used: units sold, inventory, units ordered, demand forecast, pricing, discount, weather, promotion, seasonality.",
            f"Volume: {total_units:,.0f} total units sold; average {avg_daily_units:,.0f} units/day.",
        ],
        image_path=monthly_path,
        caption="Figure: Monthly demand trend used to understand growth and seasonality.",
    )

    add_bullets_slide(
        prs,
        "3. Data Quality & Preparation",
        [
            "Column normalization maps real-world aliases to required schema (date, category/product, units sold).",
            "Date parsing supports mixed formats and protects against malformed historical values.",
            "Missing daily points are handled with resampling + forward/backward filling for robust model input.",
            "Category-level daily series are aggregated before feature generation and training.",
        ],
    )

    add_bullets_slide(
        prs,
        "4. Feature Engineering",
        [
            "Lag features: lag_1, lag_7, lag_14, lag_28 to capture immediate and medium-term demand memory.",
            "Rolling statistics: rolling mean and standard deviation for short-window trend tracking.",
            "Calendar features: day-of-week, month, weekend flags, and sinusoidal seasonal encodings.",
            "Business context inputs: inventory, pricing, discounts, competitor pricing, and promotions.",
        ],
    )

    add_bullets_slide(
        prs,
        "5. Forecasting Core",
        [
            "Current model: LinearRegression with engineered temporal + contextual features.",
            "Model supports adaptive holiday-aware dynamics through configured fixed and special holiday maps.",
            f"Historical forecast-vs-actual mean absolute gap in dataset: {abs_gap:.2f} units.",
            "Output: next-day to multi-day demand forecasts consumable by dashboard and export workflows.",
        ],
    )

    add_bullets_slide(
        prs,
        "6. System Architecture",
        [
            "Backend: FastAPI service with endpoints for forecasts, category listing, history, retraining, and admin tasks.",
            "Persistence: SQLite databases for auth/session management and forecast history.",
            "Frontend: HTML/CSS/JS dashboard pages for forecast input, result visualization, and admin operations.",
            "Upload pipeline: staged datasets with approve/reject flow for controlled replacement of production data.",
        ],
        image_path=arch_path,
        caption="Figure: High-level implementation architecture from data to decision UI.",
    )

    add_bullets_slide(
        prs,
        "7. Dashboard & Access Control",
        [
            "Role model implemented: admin, inventory_manager, and viewer with feature gating.",
            "Admin capabilities include user/dataset actions and model retraining controls.",
            "Inventory manager can run and view forecasts, with restricted admin-only operations.",
            "Viewer role is read-only for forecast consumption and monitoring.",
        ],
    )

    add_bullets_slide(
        prs,
        "8. Analytical Insights So Far",
        [
            "Demand concentration is category-skewed, indicating priority SKUs/categories for planning optimization.",
            "Regional split is uneven, so location-sensitive replenishment policies are justified.",
            f"Average inventory level is {avg_inventory:.1f}; low-stock instances (<80 units) are {low_stock_pct:.1f}%.",
            "Trend analysis suggests combining seasonal behavior with operational constraints for better procurement plans.",
        ],
        image_path=category_path,
        caption="Figure: Top category contribution by units sold.",
    )

    add_bullets_slide(
        prs,
        "9. Regional Sales Distribution",
        [
            "Region-wise demand shares are materially different and should influence stock allocation strategy.",
            "Single global reorder thresholds can be suboptimal when demand variance differs by geography.",
            "Forecast service can be extended to region-store granularity for more precise replenishment.",
            "This is a strong base for future multi-echelon inventory optimization.",
        ],
        image_path=loc_path,
        caption="Figure: Sales share across North, South, East, and West regions.",
    )

    add_bullets_slide(
        prs,
        "10. Current Limitations",
        [
            "Current baseline model is deterministic; confidence intervals are not yet surfaced in UI.",
            "Automated retraining schedules and drift alerts can be strengthened for production governance.",
            "Evaluation tracking should be expanded with a dedicated metrics dashboard over time windows.",
            "Supply constraints and lead-time uncertainty are not yet integrated into optimization logic.",
        ],
    )

    add_bullets_slide(
        prs,
        "11. Next Development Steps",
        [
            "Add probabilistic forecasting (P10/P50/P90) for risk-aware planning decisions.",
            "Introduce auto-retraining pipeline with model registry and rollback controls.",
            "Implement reorder recommendation engine combining forecast, lead time, safety stock, and service levels.",
            "Enhance admin observability with forecast error trends by category and region.",
        ],
    )

    add_bullets_slide(
        prs,
        "12. Conclusion",
        [
            "The project now has a complete functional chain: data ingestion -> feature engineering -> forecast API -> dashboard.",
            "Core foundation is stable enough for iterative optimization and production hardening.",
            "Business value path is clear: better availability, lower excess inventory, and faster planning cycles.",
            "This creates a scalable platform for AI-assisted retail supply-chain management.",
        ],
    )

    prs.save(OUT_PPT)
    print(f"Created: {OUT_PPT}")


if __name__ == "__main__":
    main()
